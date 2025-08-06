#!/usr/bin/env python3
"""
Document Processing Utility

This module handles intelligent document processing for the test generation pipeline:
1. Load documents from input directories
2. Estimate token counts and manage token limits
3. Create structured feature info with AI assistance
4. Generate and save feature input JSON files
"""

import os
import json
import tiktoken
import warnings
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any
from dataclasses import dataclass

# Suppress warnings from transformers and other libraries
warnings.filterwarnings("ignore", category=FutureWarning)
warnings.filterwarnings("ignore", message=".*attention_mask.*")
warnings.filterwarnings("ignore", message=".*_reorder_cache.*")
warnings.filterwarnings("ignore", message=".*pad_token_id.*")
warnings.filterwarnings("ignore", message=".*eos_token.*")
warnings.filterwarnings("ignore", message=".*bos_token.*")

# Also suppress at the logging level
import logging
logging.getLogger("transformers").setLevel(logging.ERROR)
logging.getLogger("torch").setLevel(logging.ERROR)

from llama_index.core import SimpleDirectoryReader
from openai import OpenAI
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

@dataclass
class DocumentInfo:
    """Information about a processed document"""
    file_name: str
    content: str
    token_count: int
    file_path: str

class DocumentProcessor:
    """
    Handles document processing with intelligent token management and AI-assisted feature info creation.
    """

    # GPT-4o token limits (leaving some buffer for system prompts and responses)
    MAX_CONTEXT_TOKENS = 120000  # GPT-4o has 128k context window
    SAFE_CONTEXT_TOKENS = 100000  # Use 100k to leave buffer for responses

    def __init__(self, api_key: Optional[str] = None, base_url: Optional[str] = None, verbose: bool = False):
        """
        Initialize the document processor.

        Args:
            api_key: OpenAI API key
            base_url: Custom base URL for OpenAI API
            verbose: Enable verbose logging
        """
        self.api_key = api_key or os.getenv("OPENAI_API_KEY")
        self.base_url = base_url or "https://apis-internal.intel.com/generativeaiinference/v4"
        self.verbose = verbose

        # Initialize OpenAI client
        if self.api_key:
            self.client = OpenAI(api_key=self.api_key, base_url=self.base_url)
        else:
            print("Warning: No OpenAI API key provided. AI features will not be available.")
            self.client = None

        # Initialize tokenizer for GPT-4o
        try:
            self.tokenizer = tiktoken.encoding_for_model("gpt-4o")
        except:
            # Fallback to cl100k_base if gpt-4o encoding is not available
            self.tokenizer = tiktoken.get_encoding("cl100k_base")

    def count_tokens(self, text: str) -> int:
        """
        Count tokens in a text string.

        Args:
            text: Text to count tokens for

        Returns:
            Number of tokens
        """
        try:
            return len(self.tokenizer.encode(text))
        except Exception as e:
            # Fallback: rough estimation (4 chars per token)
            if self.verbose:
                print(f"Warning: Token counting failed, using estimation: {e}")
            return len(text) // 4

    def load_documents_from_directory(self, docs_dir) -> List[DocumentInfo]:
        docs_path = Path(docs_dir) if isinstance(docs_dir, str) else docs_dir

        if not docs_path.exists():
            print(f"Warning: Directory {docs_path} does not exist")
            return []

        print(f"Loading documents from: {docs_path}")
        doc_infos = []
        total_tokens = 0

        all_files = []
        for file_path in docs_path.rglob('*'):
            if file_path.is_file() and not file_path.name.startswith('.'):
                all_files.append(file_path)

        if not all_files:
            print(f"No files found in {docs_path}")
            return []

        print(f"Found {len(all_files)} files to process")

        for file_path in all_files:
            try:
                # First try with SimpleDirectoryReader (supports many formats)
                content = self._load_with_simple_reader(file_path)

                if not content:
                    # Fallback to UTF-8 text reading
                    content = self._load_with_utf8_fallback(file_path)

                if content and content.strip():
                    token_count = self.count_tokens(content)

                    doc_info = DocumentInfo(
                        file_name=file_path.name,
                        content=content.strip(),
                        token_count=token_count,
                        file_path=str(file_path)
                    )
                    doc_infos.append(doc_info)
                    total_tokens += token_count

                    if self.verbose:
                        print(f"  {file_path.name}: {token_count:,} tokens")
                else:
                    if self.verbose:
                        print(f"  {file_path.name}: No content extracted")

            except Exception as e:
                print(f"  Failed to load {file_path.name}: {e}")
                continue

        print(f"Successfully loaded {len(doc_infos)} documents")
        print(f"Total content: {total_tokens:,} tokens")
        return doc_infos

    def _load_with_simple_reader(self, file_path: Path) -> str:
        try:
            problematic_extensions = {'.pptx', '.ppt'}

            if file_path.suffix.lower() in problematic_extensions:
                if self.verbose:
                    print(f"Skipping SimpleDirectoryReader for {file_path.name} (known to cause issues)")
                return ""

            if self.verbose:
                print(f"Trying SimpleDirectoryReader for {file_path.name}")

            reader = SimpleDirectoryReader(input_files=[str(file_path)])
            documents = reader.load_data()

            if documents and documents[0].text.strip():
                if self.verbose:
                    print(f"SimpleDirectoryReader succeeded for {file_path.name}")
                return documents[0].text.strip()

        except Exception as e:
            error_str = str(e).lower()
            if any(keyword in error_str for keyword in ['attention_mask', '_reorder_cache', 'transformers', 'gpt2']):
                if self.verbose:
                    print(f"    SimpleDirectoryReader failed for {file_path.name}: Transformer model issue")
            else:
                if self.verbose:
                    print(f"    SimpleDirectoryReader failed for {file_path.name}: {e}")

        return ""

    def _load_with_utf8_fallback(self, file_path: Path) -> str:
        try:
            if file_path.suffix.lower() in ['.pptx', '.ppt']:
                return self._extract_pptx_content(file_path)

            binary_extensions = {'.exe', '.bin', '.dll', '.so', '.dylib', '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.ico', '.mp3', '.mp4', '.avi', '.mov', '.zip', '.tar', '.gz', '.rar'}

            if file_path.suffix.lower() in binary_extensions:
                if self.verbose:
                    print(f"    Skipping binary file: {file_path.name}")
                return ""

            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()

            if content.strip():
                if self.verbose:
                    print(f"Loaded as UTF-8 text: {file_path.name}")
                return content.strip()

        except Exception as e:
            if self.verbose:
                print(f"UTF-8 fallback failed for {file_path.name}: {e}")

        return ""

    def _extract_pptx_content(self, pptx_file: Path) -> str:
        try:
            from pptx import Presentation

            if self.verbose:
                print(f"    Trying python-pptx for {pptx_file.name}")

            prs = Presentation(str(pptx_file))
            text_content = []

            for i, slide in enumerate(prs.slides, 1):
                text_content.append(f"=== Slide {i} ===")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text.strip())

                text_content.append("")

            result = "\n".join(text_content)
            if result.strip():
                if self.verbose:
                    print(f"Extracted PPTX content from {pptx_file.name}")
                return result.strip()

        except ImportError:
            if self.verbose:
                print(f"python-pptx not available for {pptx_file.name}")
        except Exception as e:
            if self.verbose:
                print(f"    PPTX extraction failed for {pptx_file.name}: {e}")

        return ""

    def prepare_content(self, doc_infos: List[DocumentInfo]) -> str:
        if not doc_infos:
            return ""

        total_tokens = sum(doc.token_count for doc in doc_infos)

        print(f"Preparing content for processing:")
        print(f"  Total tokens: {total_tokens:,}")
        print(f"  Safe limit: {self.SAFE_CONTEXT_TOKENS:,}")

        if total_tokens <= self.SAFE_CONTEXT_TOKENS:
            print("  Strategy: Using full document content")
            return self._combine_full_documents(doc_infos)
        else:
            print("  Strategy: Summarizing documents to fit token limits")
            return self._summarize_documents(doc_infos)

    def _combine_full_documents(self, doc_infos: List[DocumentInfo]) -> str:
        combined_text = []
        for doc in doc_infos:
            combined_text.append(f"=== Document: {doc.file_name} ===")
            combined_text.append(doc.content)
            combined_text.append("")
        return "\n".join(combined_text)

    def _summarize_documents(self, doc_infos: List[DocumentInfo]) -> str:
        if not self.client:
            print("Warning: No OpenAI client available for summarization. Using truncated content.")
            return self._truncate_documents(doc_infos)

        print("  Summarizing documents using AI...")
        summarized_docs = []

        for doc_info in doc_infos:
            try:
                summary = self._summarize_single_document(doc_info)
                if summary:
                    summarized_docs.append(f"=== {doc_info.file_name} (Summary) ===")
                    summarized_docs.append(summary)
                    summarized_docs.append("")

            except Exception as e:
                print(f"Warning: Failed to summarize {doc_info.file_name}: {e}")
                truncated = doc_info.content[:2000] + "..." if len(doc_info.content) > 2000 else doc_info.content
                summarized_docs.append(f"=== {doc_info.file_name} (Truncated) ===")
                summarized_docs.append(truncated)
                summarized_docs.append("")

        return "\n".join(summarized_docs)

    def _summarize_single_document(self, doc_info: DocumentInfo) -> str:
        prompt = f"""
        Please create a comprehensive summary of the following document that preserves all important technical details for test generation purposes.

        Document: {doc_info.file_name}

        Focus on:
        1. Main features, APIs, or components described
        2. Technical specifications and requirements
        3. Test scenarios, edge cases, or validation requirements mentioned
        4. Performance requirements or constraints
        5. Device support (CPU, GPU, etc.) and compatibility information
        6. Any specific testing instructions or requirements

        Document Content:
        {doc_info.content}

        Provide a detailed summary that captures all information relevant for generating comprehensive tests:
        """

        try:
            response = self.client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": "You are an expert technical writer who creates detailed summaries preserving all information needed for test generation."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.1,
                max_completion_tokens=2000
            )

            return response.choices[0].message.content.strip()

        except Exception as e:
            print(f"Error summarizing document {doc_info.file_name}: {e}")
            raise

    def _truncate_documents(self, doc_infos: List[DocumentInfo]) -> str:
        print("  Using truncation fallback method")
        combined_content = []
        remaining_tokens = self.SAFE_CONTEXT_TOKENS

        for doc_info in doc_infos:
            if remaining_tokens <= 1000:
                break

            if doc_info.token_count <= remaining_tokens:
                combined_content.append(f"=== {doc_info.file_name} ===")
                combined_content.append(doc_info.content)
                combined_content.append("")
                remaining_tokens -= doc_info.token_count
            else:
                chars_to_include = int((remaining_tokens / doc_info.token_count) * len(doc_info.content) * 0.8)
                truncated_content = doc_info.content[:chars_to_include] + "..."

                combined_content.append(f"=== {doc_info.file_name} (Truncated) ===")
                combined_content.append(truncated_content)
                combined_content.append("")
                break

        return "\n".join(combined_content)

    def create_feature_info_with_ai(self, content: str, input_dir_name: str) -> Dict[str, Any]:
        if not self.client:
            print("Warning: No OpenAI client available. Creating basic feature info.")
            return {
                "name": input_dir_name.replace('_', ' ').replace('-', ' ').title(),
                "description": content[:5000] + "..." if len(content) > 5000 else content
            }

        print("Creating structured feature info using AI...")

        prompt = f"""
        You are an expert at analyzing technical documentation and creating structured feature information for test generation.

        Analyze the following documentation and create a comprehensive feature description that will be used to generate thorough test plans.

        Input Directory: {input_dir_name}

        Documentation Content:
        {content}

        Please create a JSON response with the following structure:
        {{
            "name": "Descriptive Feature Name",
            "description": "Comprehensive description including all technical details, APIs, functions, test requirements, performance criteria, device support, and any specific testing instructions mentioned in the documentation."
        }}

        Instructions:
        1. Extract a clear, descriptive name for the main feature/component/API
        2. Create a comprehensive description that includes:
           - What the feature/API does and its purpose
           - All mentioned functions, methods, or APIs that need testing
           - Specific test requirements (functional, performance, edge cases, etc.)
           - Technical specifications and constraints
           - Device support requirements (CPU, GPU, HPU, etc.)
           - Data types, tensor shapes, or other technical parameters
           - Performance benchmarks or requirements
           - Any specific testing scenarios mentioned
           - Error handling and edge cases to test

        Return only the JSON object.
        """

        try:
            response = self.client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": "You are an expert at creating comprehensive feature descriptions for test generation. Always return valid JSON."},
                    {"role": "user", "content": prompt}
                ],
                response_format={"type": "json_object"},
                temperature=0.1,
                max_completion_tokens=3000
            )

            feature_info = json.loads(response.choices[0].message.content)

            if "name" in feature_info and "description" in feature_info:
                print(f"Successfully created feature info for: {feature_info['name']}")
                return feature_info
            else:
                print("Warning: AI response missing required fields. Using fallback.")
                raise ValueError("Invalid response structure")

        except Exception as e:
            print(f"Error creating feature info with AI: {e}")
            return {
                "name": input_dir_name.replace('_', ' ').replace('-', ' ').title(),
                "description": content[:10000] + "..." if len(content) > 10000 else content
            }

    def process_input_directory(self, input_dir: str, output_dir: str) -> Tuple[bool, str, Dict[str, Any]]:
        try:
            input_path = Path(input_dir)
            output_path = Path(output_dir)
            docs_dir = input_path / "docs"

            # Ensure output directory exists
            output_path.mkdir(parents=True, exist_ok=True)

            # Check if docs directory exists
            if not docs_dir.exists():
                print(f"Error: docs directory not found in {input_dir}")
                return False, "", {}
            print("=" * 60)

            # Load documents from docs directory
            doc_infos = self.load_documents_from_directory(docs_dir)

            if not doc_infos:
                print("No documents found to process")
                return False, "", {}

            # Prepare content for AI processing
            prepared_content = self.prepare_content(doc_infos)

            if not prepared_content:
                print("No meaningful content found in documents")
                return False, "", {}

            # Create feature info using AI
            input_dir_name = input_path.name
            feature_info = self.create_feature_info_with_ai(prepared_content, input_dir_name)

            # Create filename for feature input JSON
            feature_name = feature_info.get("name", input_dir_name)
            safe_name = "".join(c for c in feature_name if c.isalnum() or c in (' ', '-', '_')).strip()
            safe_name = safe_name.replace(' ', '_').lower()

            json_filename = f"{safe_name}_feature_input.json"
            json_file_path = output_path / json_filename

            # Save feature info to JSON file
            with open(json_file_path, 'w', encoding='utf-8') as f:
                json.dump(feature_info, f, indent=2, ensure_ascii=False)

            print(f"\nFeature Input JSON Created:")
            print(f"  File: {json_file_path}")
            print(f"  Feature: {feature_info['name']}")
            print(f"  Description length: {len(feature_info['description']):,} characters")

            return True, str(json_file_path), feature_info

        except Exception as e:
            print(f"Error processing input directory: {e}")
            if self.verbose:
                import traceback
                traceback.print_exc()
            return False, "", {}
