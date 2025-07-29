#!/usr/bin/env python3
"""
Multi-format Input Converter

This module converts various input file formats (PPT, DOCX, XLSX, PDF, TXT)
to the standardized JSON format used by the test generation pipeline.
"""

import os
import json
import mimetypes
from pathlib import Path
from typing import Dict, Any, Optional, List, Tuple
import tempfile
import logging

# Document processing imports
from docx import Document
import pandas as pd
from openpyxl import load_workbook

# AI processing
from openai import OpenAI
from dotenv import load_dotenv

# Optional imports for additional formats
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not available. PowerPoint files will not be supported.")

try:
    import PyPDF2
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    print("Warning: PyPDF2 not available. PDF files will not be supported.")

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class InputConverter:
    """
    Converts various input formats to standardized JSON format for test generation.
    """

    SUPPORTED_FORMATS = {
        '.json': 'JSON',
        '.txt': 'Text',
        '.docx': 'Word Document',
        '.xlsx': 'Excel Spreadsheet',
        '.xls': 'Excel Spreadsheet (Legacy)',
    }

    # Add optional formats if libraries are available
    if PPTX_AVAILABLE:
        SUPPORTED_FORMATS['.pptx'] = 'PowerPoint Presentation'

    if PDF_AVAILABLE:
        SUPPORTED_FORMATS['.pdf'] = 'PDF Document'

    def __init__(self, api_key: Optional[str] = None, base_url: Optional[str] = None):
        """
        Initialize the input converter.

        Args:
            api_key: OpenAI API key for AI conversion
            base_url: Custom base URL for OpenAI API (e.g., Intel internal)
        """
        self.api_key = api_key or os.getenv("OPENAI_API_KEY")
        self.base_url = base_url or "https://apis-internal.intel.com/generativeaiinference/v4"

        if self.api_key:
            self.client = OpenAI(api_key=self.api_key, base_url=self.base_url)
        else:
            logger.warning("No OpenAI API key provided. AI conversion will not be available.")
            self.client = None

    def get_supported_formats(self) -> Dict[str, str]:
        """Get dictionary of supported file formats."""
        return self.SUPPORTED_FORMATS.copy()

    def detect_file_format(self, file_path: str) -> str:
        """
        Args: file_path: Path to the input file
        Returns:
            File extension (e.g., '.json', '.docx')
        """
        return Path(file_path).suffix.lower()

    def is_supported_format(self, file_path: str) -> bool:
        """
        Check if the file format is supported.
        Args: file_path: Path to the input file
        Returns:
            True if supported, False otherwise
        """
        file_format = self.detect_file_format(file_path)
        return file_format in self.SUPPORTED_FORMATS

    def extract_text_from_file(self, file_path: str) -> Tuple[bool, str, str]:
        """
        Extract text content from various file formats.
        Args:
            file_path: Path to the input file
        Returns:
            Tuple of (success, extracted_text, error_message)
        """
        if not os.path.exists(file_path):
            return False, "", f"File not found: {file_path}"

        file_format = self.detect_file_format(file_path)

        if not self.is_supported_format(file_path):
            return False, "", f"Unsupported file format: {file_format}"

        try:
            if file_format == '.json' or file_format == '.txt':
                return self._extract_content(file_path)
            elif file_format == '.docx':
                return self._extract_from_docx(file_path)
            elif file_format in ['.xlsx', '.xls']:
                return self._extract_from_excel(file_path)
            elif file_format == '.pptx' and PPTX_AVAILABLE:
                return self._extract_from_pptx(file_path)
            elif file_format == '.pdf' and PDF_AVAILABLE:
                return self._extract_from_pdf(file_path)
            else:
                return False, "", f"No handler available for format: {file_format}"

        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            return False, "", f"Error processing file: {str(e)}"

    def _extract_content(self, file_path: str) -> Tuple[bool, str, str]:
        """Extract content from JSON or txt file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            return True, content, ""
        except Exception as e:
            return False, "", str(e)

    def _extract_from_docx(self, file_path: str) -> Tuple[bool, str, str]:
        """Extract text from Word document."""
        try:
            doc = Document(file_path)
            text_content = []

            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text.strip())

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(" | ".join(row_text))

            return True, "\n".join(text_content), ""
        except Exception as e:
            return False, "", str(e)

    def _extract_from_excel(self, file_path: str) -> Tuple[bool, str, str]:
        """Extract text from Excel spreadsheet."""
        try:
            # Try with pandas first (supports both .xlsx and .xls)
            excel_file = pd.ExcelFile(file_path)
            text_content = []

            for sheet_name in excel_file.sheet_names:
                text_content.append(f"=== Sheet: {sheet_name} ===")
                df = pd.read_excel(file_path, sheet_name=sheet_name)

                # Convert DataFrame to string representation
                sheet_text = df.to_string(index=False, na_rep='')
                text_content.append(sheet_text)
                text_content.append("")  # Add blank line between sheets

            return True, "\n".join(text_content), ""
        except Exception as e:
            return False, "", str(e)

    def _extract_from_pptx(self, file_path: str) -> Tuple[bool, str, str]:
        """Extract text from PowerPoint presentation."""
        if not PPTX_AVAILABLE:
            return False, "", "PowerPoint support not available (python-pptx not installed)"

        try:
            prs = Presentation(file_path)
            text_content = []

            for i, slide in enumerate(prs.slides, 1):
                text_content.append(f"=== Slide {i} ===")

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text.strip())

                text_content.append("")  # Add blank line between slides

            return True, "\n".join(text_content), ""
        except Exception as e:
            return False, "", str(e)

    def _extract_from_pdf(self, file_path: str) -> Tuple[bool, str, str]:
        """Extract text from PDF document."""
        if not PDF_AVAILABLE:
            return False, "", "PDF support not available (PyPDF2 not installed)"

        try:
            text_content = []

            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)

                for i, page in enumerate(pdf_reader.pages, 1):
                    text_content.append(f"=== Page {i} ===")
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_content.append(page_text.strip())
                    text_content.append("")  # Add blank line between pages

            return True, "\n".join(text_content), ""
        except Exception as e:
            return False, "", str(e)

    def convert_to_json(self, file_path: str, output_file: Optional[str] = None) -> Tuple[bool, str, Dict[str, Any]]:
        """
        Convert input file to standardized JSON format.

        Args:
            file_path: Path to input file
            output_file: Optional path to save converted JSON

        Returns:
            Tuple of (success, output_path, json_content)
        """
        file_format = self.detect_file_format(file_path)
        input_file_name = Path(file_path).name
        format_description = self.SUPPORTED_FORMATS.get(file_format, "Unknown")

        # Print input format information
        print(f"\nInput File Analysis:")
        print(f"   File: {input_file_name}")
        print(f"   Format: {format_description} ({file_format})")

        # If already JSON, validate and return as-is
        if file_format == '.json':
            print(f"   Status: Validating existing JSON structure...")
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    json_content = json.load(f)

                # Validate that it has required structure
                if self._validate_json_structure(json_content):
                    print(f"   JSON file is already in correct format")
                    # Even for valid JSON, save a copy if output_file is specified
                    if output_file and output_file != file_path:
                        try:
                            with open(output_file, 'w', encoding='utf-8') as f:
                                json.dump(json_content, f, indent=2)
                            print(f"   Copied to: {output_file}")
                            return True, output_file, json_content
                        except Exception as e:
                            logger.error(f"Error copying JSON file: {e}")
                    return True, file_path, json_content
                else:
                    print(f"   JSON file doesn't match expected structure, converting via AI...")
            except Exception as e:
                print(f"   Error reading JSON file: {e}")
                return False, "", {}

        # Extract text content
        print(f"   Status: Extracting content from {format_description}...")
        success, text_content, error = self.extract_text_from_file(file_path)
        if not success:
            print(f"   Failed to extract content: {error}")
            return False, "", {}

        print(f"   Successfully extracted {len(text_content)} characters")

        # Convert to JSON using AI
        print(f"\nAI Conversion Process:")
        print(f"   Converting {format_description} content to standardized JSON format...")

        if not self.client:
            print(f"   No OpenAI client available for AI conversion")
            logger.error("No OpenAI client available for AI conversion")
            return False, "", {}

        json_content = self._convert_text_to_json_with_ai(text_content, file_path)
        if not json_content:
            print(f"   AI conversion failed")
            return False, "", {}

        print(f"   Successfully converted to JSON format")

        # Determine output file path
        if not output_file:
            # Create output file name based on input file
            input_stem = Path(file_path).stem
            output_file = str(Path(file_path).parent / f"{input_stem}_converted.json")

        # Always save to disk
        try:
            print(f"\nSaving Converted JSON:")
            with open(output_file, 'w', encoding='utf-8') as f:
                json.dump(json_content, f, indent=2)
            print(f"   Successfully saved to: {output_file}")
            print(f"   JSON contains: name='{json_content.get('name', 'N/A')}', description length={len(json_content.get('description', ''))}")
            return True, output_file, json_content
        except Exception as e:
            print(f"   Error saving JSON file: {e}")
            logger.error(f"Error saving JSON file: {e}")
            return False, "", {}

    def _validate_json_structure(self, json_content: Dict[str, Any]) -> bool:
        required_fields = ['name', 'description']
        return all(field in json_content for field in required_fields)

    def _convert_text_to_json_with_ai(self, text_content: str, original_file: str) -> Optional[Dict[str, Any]]:
        """
        Convert extracted text to JSON using AI.

        Args:
            text_content: Extracted text content
            original_file: Path to original file for context

        Returns:
            JSON content or None if conversion failed
        """
        try:
            file_name = Path(original_file).name

            prompt = f"""
            You are an expert at converting feature documentation into standardized JSON format for test generation.

            Please convert the following content from "{file_name}" into the exact JSON format specified below:

            REQUIRED JSON FORMAT:
            {{
            "name": "Feature Name",
            "description": "Detailed description of the feature that includes what needs to be tested. Include information about APIs, functions, methods, or features that need validation. Specify test requirements like performance testing, edge cases, functional tests, zero-size tensor tests, big tensor tests, etc. Include details about data types, tensor shapes, world size, CPU/GPU/HPU support, synchronization, scaling, and memory requirements."
            }}

            CONTENT TO CONVERT:
            {text_content}

            IMPORTANT INSTRUCTIONS:
            1. Extract the main feature/API/component name for the "name" field
            2. Create a comprehensive description that includes:
            - What the feature does
            - What APIs/functions need testing
            - Types of tests needed (performance, functional, edge cases, etc.)
            - Parameters to test (data types, tensor shapes, world size)
            - Device support requirements (CPU, GPU, HPU)
            - Performance and scaling requirements
            3. If the content mentions specific APIs, include them in the description
            4. If test requirements are mentioned, include them
            5. Make the description detailed enough for comprehensive test generation
            6. Return ONLY the JSON object, no additional text

            JSON Response:
            """

            response = self.client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": "You are an expert at converting feature documentation to JSON format for test generation. Always return valid JSON."},
                    {"role": "user", "content": prompt}
                ],
                response_format={"type": "json_object"},
                temperature=0.1,
                max_completion_tokens=2000
            )

            response_text = response.choices[0].message.content
            json_content = json.loads(response_text)

            # Validate the response
            if self._validate_json_structure(json_content):
                return json_content
            else:
                logger.error("AI response doesn't match expected JSON structure")
                return None

        except Exception as e:
            logger.error(f"Error converting text to JSON with AI: {e}")
            return None